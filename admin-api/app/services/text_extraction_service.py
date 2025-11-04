"""
Text Extraction Service for Various Document Formats
다양한 문서 형식에서 텍스트 추출 서비스
"""
from typing import BinaryIO, Optional
import os


class TextExtractionService:
    """문서에서 텍스트 추출"""

    def extract_text(self, file_obj: BinaryIO, filename: str, mime_type: str) -> str:
        """
        파일에서 텍스트 추출
        """
        file_ext = os.path.splitext(filename)[1].lower()

        try:
            if file_ext == '.txt' or mime_type == 'text/plain':
                return self._extract_from_txt(file_obj)

            elif file_ext == '.pdf' or mime_type == 'application/pdf':
                return self._extract_from_pdf(file_obj)

            elif file_ext == '.docx' or 'wordprocessingml' in mime_type:
                return self._extract_from_docx(file_obj)

            elif file_ext in ['.hwp', '.hwpx'] or 'hwp' in mime_type.lower():
                return self._extract_from_hwp(file_obj, file_ext)

            elif file_ext in ['.ppt', '.pptx'] or 'presentationml' in mime_type or 'powerpoint' in mime_type:
                return self._extract_from_ppt(file_obj, file_ext)

            else:
                # Fallback: try to read as text
                return self._extract_from_txt(file_obj)

        except Exception as e:
            print(f"Text extraction error for {filename}: {e}")
            # Return minimal text on error
            return f"Document: {filename}\n[텍스트 추출 실패: {str(e)}]"

    def _extract_from_txt(self, file_obj: BinaryIO) -> str:
        """TXT 파일 텍스트 추출"""
        content = file_obj.read()
        # Try multiple encodings
        for encoding in ['utf-8', 'cp949', 'euc-kr', 'latin-1']:
            try:
                return content.decode(encoding)
            except UnicodeDecodeError:
                continue
        return content.decode('utf-8', errors='ignore')

    def _extract_from_pdf(self, file_obj: BinaryIO) -> str:
        """PDF 파일 텍스트 추출"""
        try:
            import PyPDF2
            pdf_reader = PyPDF2.PdfReader(file_obj)
            text_parts = []
            for page in pdf_reader.pages:
                page_text = page.extract_text()
                # Remove NULL bytes that PostgreSQL cannot handle
                if page_text:
                    page_text = page_text.replace('\x00', '')
                text_parts.append(page_text)
            return '\n'.join(text_parts)
        except ImportError:
            return "[PDF 텍스트 추출 라이브러리 미설치: PyPDF2]"
        except Exception as e:
            return f"[PDF 추출 실패: {str(e)}]"

    def _extract_from_docx(self, file_obj: BinaryIO) -> str:
        """DOCX 파일 텍스트 추출"""
        try:
            from docx import Document
            doc = Document(file_obj)
            text_parts = []
            for paragraph in doc.paragraphs:
                # Remove NULL bytes that PostgreSQL cannot handle
                text = paragraph.text.replace('\x00', '') if paragraph.text else ''
                text_parts.append(text)
            return '\n'.join(text_parts)
        except ImportError:
            return "[DOCX 텍스트 추출 라이브러리 미설치: python-docx]"
        except Exception as e:
            return f"[DOCX 추출 실패: {str(e)}]"

    def _extract_from_hwp(self, file_obj: BinaryIO, file_ext: str) -> str:
        """HWP/HWPX 파일 텍스트 추출"""
        try:
            # HWPX는 ZIP 기반 XML 형식
            if file_ext == '.hwpx':
                import zipfile
                import xml.etree.ElementTree as ET

                with zipfile.ZipFile(file_obj) as zf:
                    # Contents/section0.xml에서 텍스트 추출
                    text_parts = []
                    for name in zf.namelist():
                        if 'section' in name and name.endswith('.xml'):
                            content = zf.read(name)
                            root = ET.fromstring(content)
                            # Extract text from all text nodes
                            for text_elem in root.iter():
                                if text_elem.text:
                                    text_parts.append(text_elem.text.strip())
                    return '\n'.join(filter(None, text_parts))

            else:  # .hwp
                # HWP는 OLE 기반, olefile 사용
                import olefile
                ole = olefile.OleFileIO(file_obj)

                # Try to extract from BodyText sections
                text_parts = []
                for entry in ole.listdir():
                    if 'BodyText' in entry or 'Section' in entry:
                        try:
                            stream = ole.openstream(entry)
                            data = stream.read()
                            # Try to decode as Korean text
                            for encoding in ['utf-16le', 'cp949', 'utf-8']:
                                try:
                                    text = data.decode(encoding, errors='ignore')
                                    if text.strip():
                                        text_parts.append(text)
                                    break
                                except:
                                    continue
                        except:
                            continue

                if text_parts:
                    return '\n'.join(text_parts)
                else:
                    return "[HWP 텍스트 추출 부분 성공: 구조 분석 필요]"

        except ImportError as e:
            return f"[HWP 텍스트 추출 라이브러리 미설치: {str(e)}]"
        except Exception as e:
            return f"[HWP 추출 실패: {str(e)}]"

    def _extract_from_ppt(self, file_obj: BinaryIO, file_ext: str) -> str:
        """PPT/PPTX 파일 텍스트 추출"""
        try:
            if file_ext == '.pptx':
                from pptx import Presentation
                prs = Presentation(file_obj)
                text_parts = []

                for slide_num, slide in enumerate(prs.slides, 1):
                    text_parts.append(f"\n=== Slide {slide_num} ===")
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_parts.append(shape.text)

                return '\n'.join(text_parts)

            else:  # .ppt (OLE format)
                # PPT는 복잡한 바이너리 형식, 기본 추출만
                return "[PPT 파일: PPTX 형식으로 변환 후 업로드 권장]"

        except ImportError:
            return "[PPT 텍스트 추출 라이브러리 미설치: python-pptx]"
        except Exception as e:
            return f"[PPT 추출 실패: {str(e)}]"


# Singleton instance
text_extraction_service = TextExtractionService()
